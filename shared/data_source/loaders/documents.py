"""
文档加载器实现

支持多种文档格式的加载和内容提取
"""
from typing import Optional
import io

from .base import BaseLoader
from ..schemas.models import RawDocument, UnifiedDocument, SourceType


class ExcelLoader(BaseLoader):
    """Excel 文件加载器"""

    def __init__(self):
        super().__init__(SourceType.EXCEL)

    async def load(self, raw_doc: RawDocument) -> Optional[UnifiedDocument]:
        """加载 Excel 文件"""
        try:
            # 需要安装: pip install openpyxl pandas
            import pandas as pd

            # 读取 Excel 文件
            if isinstance(raw_doc.raw_data, bytes):
                excel_file = io.BytesIO(raw_doc.raw_data)
            else:
                excel_file = raw_doc.source_id

            # 读取所有工作表
            excel_data = pd.read_excel(excel_file, sheet_name=None)

            # 将所有工作表转换为文本
            content_parts = []
            for sheet_name, df in excel_data.items():
                content_parts.append(f"## 工作表: {sheet_name}\n")
                # 转换为 markdown 表格
                content_parts.append(df.to_markdown(index=False))
                content_parts.append("\n")

            content = "\n".join(content_parts)

            return self._create_base_document(
                raw_doc=raw_doc,
                content=content,
                title=raw_doc.metadata.get('filename', 'Excel文档')
            )

        except Exception as e:
            print(f"加载 Excel 文件失败: {e}")
            return None


class WordLoader(BaseLoader):
    """Word 文件加载器"""

    def __init__(self):
        super().__init__(SourceType.WORD)

    async def load(self, raw_doc: RawDocument) -> Optional[UnifiedDocument]:
        """加载 Word 文件"""
        try:
            # 需要安装: pip install python-docx
            from docx import Document

            if isinstance(raw_doc.raw_data, bytes):
                doc_file = io.BytesIO(raw_doc.raw_data)
            else:
                doc_file = raw_doc.source_id

            doc = Document(doc_file)

            # 提取所有段落
            content_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)

            # 提取表格
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                if table_text:
                    content_parts.append("\n表格:\n" + "\n".join(table_text))

            content = "\n\n".join(content_parts)

            return self._create_base_document(
                raw_doc=raw_doc,
                content=content,
                title=raw_doc.metadata.get('filename', 'Word文档')
            )

        except Exception as e:
            print(f"加载 Word 文件失败: {e}")
            return None


class PowerPointLoader(BaseLoader):
    """PowerPoint 文件加载器"""

    def __init__(self):
        super().__init__(SourceType.POWERPOINT)

    async def load(self, raw_doc: RawDocument) -> Optional[UnifiedDocument]:
        """加载 PowerPoint 文件"""
        try:
            # 需要安装: pip install python-pptx
            from pptx import Presentation

            if isinstance(raw_doc.raw_data, bytes):
                ppt_file = io.BytesIO(raw_doc.raw_data)
            else:
                ppt_file = raw_doc.source_id

            prs = Presentation(ppt_file)

            # 提取所有幻灯片的文本
            content_parts = []
            for i, slide in enumerate(prs.slides, 1):
                content_parts.append(f"## 幻灯片 {i}\n")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)

                content_parts.append("\n")

            content = "\n".join(content_parts)

            return self._create_base_document(
                raw_doc=raw_doc,
                content=content,
                title=raw_doc.metadata.get('filename', 'PowerPoint文档')
            )

        except Exception as e:
            print(f"加载 PowerPoint 文件失败: {e}")
            return None


class PDFLoader(BaseLoader):
    """PDF 文件加载器"""

    def __init__(self):
        super().__init__(SourceType.PDF)

    async def load(self, raw_doc: RawDocument) -> Optional[UnifiedDocument]:
        """加载 PDF 文件"""
        try:
            # 需要安装: pip install pypdf
            from pypdf import PdfReader

            if isinstance(raw_doc.raw_data, bytes):
                pdf_file = io.BytesIO(raw_doc.raw_data)
            else:
                pdf_file = raw_doc.source_id

            reader = PdfReader(pdf_file)

            # 提取所有页面的文本
            content_parts = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    content_parts.append(f"## 第 {i} 页\n")
                    content_parts.append(text)
                    content_parts.append("\n")

            content = "\n".join(content_parts)

            # 提取元数据
            metadata = raw_doc.metadata.copy()
            if reader.metadata:
                metadata.update({
                    'pdf_title': reader.metadata.get('/Title'),
                    'pdf_author': reader.metadata.get('/Author'),
                    'pdf_subject': reader.metadata.get('/Subject'),
                    'pdf_pages': len(reader.pages),
                })

            return self._create_base_document(
                raw_doc=raw_doc,
                content=content,
                title=metadata.get('pdf_title') or raw_doc.metadata.get('filename', 'PDF文档'),
                author=metadata.get('pdf_author'),
                metadata=metadata
            )

        except Exception as e:
            print(f"加载 PDF 文件失败: {e}")
            return None


class ImageLoader(BaseLoader):
    """图像文件加载器（使用 OCR 提取文本）"""

    def __init__(self):
        super().__init__(SourceType.IMAGE)

    async def load(self, raw_doc: RawDocument) -> Optional[UnifiedDocument]:
        """加载图像文件"""
        try:
            # 需要安装: pip install pytesseract pillow
            # 并安装 Tesseract OCR: https://github.com/tesseract-ocr/tesseract
            from PIL import Image
            import pytesseract

            if isinstance(raw_doc.raw_data, bytes):
                image = Image.open(io.BytesIO(raw_doc.raw_data))
            else:
                image = Image.open(raw_doc.source_id)

            # 使用 OCR 提取文本
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')

            # 获取图像信息
            metadata = raw_doc.metadata.copy()
            metadata.update({
                'image_width': image.width,
                'image_height': image.height,
                'image_format': image.format,
                'image_mode': image.mode,
            })

            return self._create_base_document(
                raw_doc=raw_doc,
                content=text,
                title=raw_doc.metadata.get('filename', '图像文档'),
                metadata=metadata
            )

        except Exception as e:
            print(f"加载图像文件失败: {e}")
            print("提示: 图像 OCR 需要安装 pytesseract 和 Tesseract OCR")
            return None
